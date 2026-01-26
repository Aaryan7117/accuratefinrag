"""
PowerPoint Parser for Financial Presentations
==============================================

This module handles PowerPoint file parsing for financial presentations and reports.
Supports .pptx and .ppt formats with slide structure preservation.

Design Decisions:
- Use python-pptx for robust PowerPoint reading
- Treat each slide as a document section
- Extract text from shapes, tables, and notes
- Preserve slide titles and content hierarchy
"""

import hashlib
from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Optional, Dict, Any, Tuple
from enum import Enum

from pptx import Presentation


class ElementType(Enum):
    """Types of document elements."""
    HEADING = "heading"
    PARAGRAPH = "paragraph"
    TABLE = "table"
    LIST = "list"
    UNKNOWN = "unknown"


@dataclass
class DocumentElement:
    """
    A single element from the parsed document.

    Every element tracks its source location for citation purposes.
    """
    element_type: ElementType
    content: str
    page_number: int  # For PPT, this will be slide number
    bbox: Optional[Tuple[float, float, float, float]] = None  # Not applicable for PPT
    heading_level: Optional[int] = None
    table_id: Optional[str] = None
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class Section:
    """
    A document section containing multiple elements.

    For PowerPoint, each slide becomes a section.
    """
    section_id: str
    title: str
    level: int = 1  # All slides are top-level
    item_number: Optional[str] = None
    page_start: int = 0
    page_end: int = 0
    elements: List[DocumentElement] = field(default_factory=list)
    subsections: List["Section"] = field(default_factory=list)


@dataclass
class ParsedDocument:
    """
    Complete parsed document with structure.

    Contains all content from PowerPoint slides.
    """
    doc_id: str
    filename: str
    title: Optional[str] = None
    total_pages: int = 0  # Number of slides
    sections: List[Section] = field(default_factory=list)
    elements: List[DocumentElement] = field(default_factory=list)
    tables: List[DocumentElement] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)

    def get_all_text(self) -> str:
        """Get all text content concatenated."""
        return "\n\n".join(
            el.content for el in self.elements
            if el.element_type != ElementType.TABLE
        )


class PowerPointParser:
    """
    PowerPoint parser for financial presentations.

    Extraction Strategy:
    1. Read all slides from the presentation
    2. Extract text from shapes, tables, and notes
    3. Create sections for each slide
    4. Preserve slide titles and structure
    """

    def __init__(self):
        """Initialize parser."""
        pass

    def parse(self, ppt_path: str | Path) -> ParsedDocument:
        """
        Parse a PowerPoint document.

        Args:
            ppt_path: Path to the PowerPoint file

        Returns:
            ParsedDocument with all extracted content
        """
        ppt_path = Path(ppt_path)

        # Generate stable document ID from file hash
        doc_id = self._generate_doc_id(ppt_path)

        # Initialize document
        doc = ParsedDocument(
            doc_id=doc_id,
            filename=ppt_path.name,
            metadata={
                "source_path": str(ppt_path),
                "file_size": ppt_path.stat().st_size,
            }
        )

        # Load presentation
        try:
            presentation = Presentation(ppt_path)
        except Exception as e:
            raise ValueError(f"Failed to read PowerPoint file: {e}")

        doc.total_pages = len(presentation.slides)

        # Process each slide
        for slide_idx, slide in enumerate(presentation.slides):
            # Get slide title
            slide_title = self._extract_slide_title(slide)
            if not slide_title:
                slide_title = f"Slide {slide_idx + 1}"

            # Create section for this slide
            section = Section(
                section_id=f"slide_{slide_idx}",
                title=slide_title,
                level=1,
                page_start=slide_idx,
                page_end=slide_idx
            )

            # Extract content from slide
            slide_elements = self._extract_slide_content(slide, slide_idx)
            section.elements.extend(slide_elements)
            doc.elements.extend(slide_elements)

            # Extract tables separately
            slide_tables = self._extract_slide_tables(slide, slide_idx)
            section.elements.extend(slide_tables)
            doc.tables.extend(slide_tables)

            doc.sections.append(section)

        # Set title to first slide title or filename
        if doc.sections:
            doc.title = doc.sections[0].title
        else:
            doc.title = ppt_path.stem

        return doc

    def _generate_doc_id(self, ppt_path: Path) -> str:
        """
        Generate stable document ID from file content hash.
        """
        hasher = hashlib.sha256()
        with open(ppt_path, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                hasher.update(chunk)
        return hasher.hexdigest()[:16]

    def _extract_slide_title(self, slide) -> str:
        """
        Extract the title from a slide.
        """
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Check if this is a title shape (usually the first text shape)
                return shape.text.strip()
        return ""

    def _extract_slide_content(self, slide, slide_idx: int) -> List[DocumentElement]:
        """
        Extract all text content from a slide.
        """
        elements = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Skip if this is likely a title (handled separately)
                if len(shape.text.strip()) < 100 and shape.text.isupper():
                    continue

                element = DocumentElement(
                    element_type=ElementType.PARAGRAPH,
                    content=shape.text.strip(),
                    page_number=slide_idx,
                    metadata={
                        "shape_type": str(type(shape).__name__),
                        "has_text": True
                    }
                )
                elements.append(element)

        # Extract notes if available
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                element = DocumentElement(
                    element_type=ElementType.PARAGRAPH,
                    content=f"Notes: {notes_text}",
                    page_number=slide_idx,
                    metadata={"content_type": "notes"}
                )
                elements.append(element)

        return elements

    def _extract_slide_tables(self, slide, slide_idx: int) -> List[DocumentElement]:
        """
        Extract tables from a slide.
        """
        tables = []
        table_counter = 0

        for shape in slide.shapes:
            if hasattr(shape, "table"):
                table_counter += 1
                table_id = f"table_{slide_idx}_{table_counter}"

                # Convert table to markdown
                table_content = self._table_to_markdown(shape.table)

                table_element = DocumentElement(
                    element_type=ElementType.TABLE,
                    content=table_content,
                    page_number=slide_idx,
                    table_id=table_id,
                    metadata={
                        "rows": len(shape.table.rows),
                        "cols": len(shape.table.columns) if shape.table.columns else 0,
                        "shape_type": "table"
                    }
                )
                tables.append(table_element)

        return tables

    def _table_to_markdown(self, table) -> str:
        """
        Convert PowerPoint table to markdown format.
        """
        if not table.rows:
            return ""

        lines = []
        for i, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                cells.append(cell_text)

            lines.append("| " + " | ".join(cells) + " |")

            # Add header separator after first row
            if i == 0:
                lines.append("|" + "|".join(["---"] * len(cells)) + "|")

        return "\n".join(lines)


# Convenience function for simple usage
def parse_ppt(ppt_path: str | Path) -> ParsedDocument:
    """Parse a PowerPoint file and return structured document."""
    parser = PowerPointParser()
    return parser.parse(ppt_path)